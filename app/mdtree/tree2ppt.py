import datetime
import os
from enum import Enum
from io import BytesIO
import requests
import json
from urllib.parse import quote_plus
from app.core.configs import settings
import markdown
from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.text.text import Font

from app.mdtree.parser import parse_string, Out, Heading
from app.mdtree.utils import get_random_theme, get_random_file, read_md_file



class Tree2PPT:
    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None

    def __init__(self, md_str1):
        self.init_pptx()
        self.init_markdown(md_str1)
        self.traverse_tree(self.tree)
        now = datetime.datetime.now().timestamp()
        path = './myppt/test' + str(now) + '.pptx'
        if not os.path.exists('./myppt'):
            os.makedirs('./myppt')
        self.prs.save(path)
        pass

    def init_pptx(self):
        prs = Presentation()
        self.theme = get_random_theme()
        self.prs = prs

    def init_markdown(self, md_str):
        self.md_str = md_str
        self.out = parse_string(md_str)
        self.tree = self.out.main

    def traverse_tree(self, heading):
        if heading is not None and (heading.source is None or heading.source == ''):
            content = ""
            if heading.children is not []:
                for child in heading.children:
                    content = content + child.text + "\n"
            MD2Slide(self.prs, self.theme, heading.text, content=content)
        elif heading is not None:
            MD2Slide(self.prs, self.theme, heading.text, content=heading.source)
        else:
            return

        # self.make_slide_demo(self.prs, heading.text, heading.source)
        if heading.children is not []:
            for child in heading.children:
                self.traverse_tree(child)

    def save_stream(self):
        stream = BytesIO()
        self.prs.save(stream)
        stream.seek(0)  # Reset the stream position to the beginning
        return stream


class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"

    pass


class MD2Slide:
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "Arial"
    font_title_size: Pt = Pt(26)
    font_content_size: Pt = Pt(18)
    font_title_color: rgb = RGBColor(51, 0, 102)
    font_content_color: rgb = RGBColor(51, 0, 102)

    def __init__(self, presentation, theme_path, title, content, *args, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[8])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()
        self.insert_image()

    def init_slide(self):
        placeholder1 = self.slide.placeholders[1]
        path = get_random_file(self.theme)
        picture = placeholder1.insert_picture(path)
        placeholder2 = self.slide.placeholders[2]
        placeholder2.element.getparent().remove(placeholder2.element)
        # 2、设置占位符宽高
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        if 'font_name' in kwargs:
            self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs:
            self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs:
            self.font_content_size = kwargs['font_content_size']
        if 'font_title_color' in kwargs:
            self.font_title_color = kwargs['font_title_color']
        if 'font_content_color' in kwargs:
            self.font_content_color = kwargs['font_content_color']

    def get_font(self, font: Font, category: str):
        font.bold = True
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.8))
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # 添加标题
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        tf = text_box_content.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        # 添加正文
        paragraph = tf.paragraphs[0]
        paragraph.text = self.content.replace("<p>", "").replace("</p>", "\n")
        self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
        # TODO 处理正文
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    def search_pexels_images(self, keyword):
        query = quote_plus(keyword.lower())
        PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
        headers = {
            'Authorization': settings.PEXELS_API
        }
        response = requests.get(PEXELS_API_URL, headers=headers)
        data = json.loads(response.text)
        if 'photos' in data:
            if len(data['photos']) > 0:
                return data['photos'][0]['src']['medium']
        return None
    def insert_image(self):
        image_url = self.search_pexels_images(self.title)
        if image_url is not None:
            # Tải ảnh từ URL
            image_data = requests.get(image_url).content
            image_stream = BytesIO(image_data)
            
            # Kích thước slide
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            
            # Kích thước ảnh (vừa phải)
            image_width = Inches(6)  # Chiều rộng ảnh
            image_height = Inches(4)  # Chiều cao ảnh
            
            # Tính toán vị trí góc dưới bên phải
            left = slide_width - image_width - Inches(0.5)  # Cách lề phải 0.5 inch
            top = slide_height - image_height - Inches(0.5)  # Cách lề dưới 0.5 inch
            
            # Thêm ảnh vào slide
            picture = self.slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
            
            # Thêm viền và đổ bóng cho ảnh (tùy chọn)
            picture.line.color.rgb = RGBColor(0, 0, 0)  # Viền màu đen
            picture.line.width = Pt(1.5)  # Độ dày viền
            picture.shadow.inherit = False  # Tắt đổ bóng mặc định (nếu cần)
            
            
    def processing_md_str(self, md_str):
        print(md_str)
        md = markdown.Markdown()
        html1 = md.convert(md_str)
        print(html1)
